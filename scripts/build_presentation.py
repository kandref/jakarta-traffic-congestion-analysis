from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORT_DIR = ROOT / "reports"
OUTPUT_PATH = REPORT_DIR / "jakarta_traffic_congestion_findings.pptx"

RAW_PATH = ROOT / "data" / "raw" / "jakarta_41_corridor_peak_speed.csv"
CLEAN_PATH = ROOT / "data" / "processed" / "jakarta_corridor_peak_speed_clean.csv"
RANKING_PATH = ROOT / "data" / "processed" / "jakarta_corridor_bottleneck_ranking.csv"
MONTHLY_PATH = ROOT / "data" / "processed" / "jakarta_corridor_monthly_speed_summary.csv"
SLOWEST_FIGURE = ROOT / "results" / "figures" / "jakarta_slowest_corridors.png"
MONTHLY_FIGURE = ROOT / "results" / "figures" / "jakarta_monthly_peak_speed.png"

NAVY = RGBColor(20, 38, 61)
BLUE = RGBColor(35, 90, 155)
RED = RGBColor(194, 59, 34)
GREEN = RGBColor(43, 128, 92)
GRAY = RGBColor(94, 105, 120)
LIGHT_GRAY = RGBColor(242, 245, 248)
WHITE = RGBColor(255, 255, 255)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.25), Inches(0.75))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(
            Inches(0.58),
            Inches(1.05),
            Inches(11.8),
            Inches(0.35),
        )
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = GRAY


def add_footer(slide, text: str = "Jakarta Traffic Congestion Analysis") -> None:
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(6.92), Inches(12.25), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.color.rgb = LIGHT_GRAY
    footer = slide.shapes.add_textbox(Inches(0.58), Inches(6.98), Inches(12), Inches(0.2))
    p = footer.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(17)
        p.font.color.rgb = NAVY
        p.space_after = Pt(8)


def add_metric_card(slide, left: float, top: float, width: float, title: str, value: str, color) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(1.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = RGBColor(220, 226, 232)
    text = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.12), Inches(width - 0.36), Inches(0.8))
    tf = text.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = value
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = color
    q = tf.add_paragraph()
    q.text = title
    q.font.size = Pt(10)
    q.font.color.rgb = GRAY


def add_table(slide, df: pd.DataFrame, left: float, top: float, width: float, height: float) -> None:
    table_shape = slide.shapes.add_table(
        rows=len(df) + 1,
        cols=len(df.columns),
        left=Inches(left),
        top=Inches(top),
        width=Inches(width),
        height=Inches(height),
    )
    table = table_shape.table
    for col_idx, column in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = column
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = WHITE

    for row_idx, row in enumerate(df.itertuples(index=False), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 else LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = NAVY
                paragraph.alignment = PP_ALIGN.LEFT if col_idx == 1 else PP_ALIGN.CENTER


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def build_deck() -> None:
    clean = pd.read_csv(CLEAN_PATH)
    ranking = pd.read_csv(RANKING_PATH)
    monthly = pd.read_csv(MONTHLY_PATH)

    overall_mean = clean["speed_kmh"].mean()
    corridor_count = clean["corridor"].nunique()
    slowest = ranking.iloc[0]
    fastest = ranking.sort_values("mean_speed_kmh", ascending=False).iloc[0]
    slowest_month = monthly.loc[monthly["mean_speed_kmh"].idxmin()]
    fastest_month = monthly.loc[monthly["mean_speed_kmh"].idxmax()]
    below_30 = int((ranking["mean_speed_kmh"] < 30).sum())

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = blank_slide(prs)
    banner = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.color.rgb = NAVY
    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.8), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "Jakarta Traffic Congestion Analysis"
    p.font.bold = True
    p.font.size = Pt(42)
    p.font.color.rgb = WHITE
    subtitle = slide.shapes.add_textbox(Inches(0.85), Inches(2.82), Inches(10.8), Inches(0.8))
    sp = subtitle.text_frame.paragraphs[0]
    sp.text = "Early findings from 2023 peak-hour corridor speed data"
    sp.font.size = Pt(22)
    sp.font.color.rgb = RGBColor(214, 225, 238)
    footer = slide.shapes.add_textbox(Inches(0.85), Inches(6.65), Inches(10), Inches(0.3))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "Data source: Satu Data Jakarta / Data.go.id"
    fp.font.size = Pt(12)
    fp.font.color.rgb = RGBColor(214, 225, 238)

    slide = blank_slide(prs)
    add_title(slide, "Executive Summary", "A first-pass diagnosis using corridor speed as congestion signal")
    add_bullets(
        slide,
        [
            f"The dataset has {len(clean):,} observations across 12 months in 2023.",
            f"The raw file contains {corridor_count} unique corridor names, although the catalog title mentions 41 corridors.",
            f"Overall mean peak-hour speed is {overall_mean:.2f} km/h.",
            f"{below_30} corridors have mean peak-hour speed below 30 km/h.",
            "The pattern points to persistent corridor-level bottlenecks, not only a one-off monthly fluctuation.",
        ],
        0.75,
        1.55,
        11.8,
        4.7,
    )
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "Dataset Snapshot", "Current analysis uses corridor speed; demand and road-capacity data are next layers")
    add_metric_card(slide, 0.75, 1.55, 2.8, "observations", f"{len(clean):,}", BLUE)
    add_metric_card(slide, 3.85, 1.55, 2.8, "months in 2023", "12", GREEN)
    add_metric_card(slide, 6.95, 1.55, 2.8, "unique corridor names", f"{corridor_count}", BLUE)
    add_metric_card(slide, 10.05, 1.55, 2.8, "mean speed", f"{overall_mean:.1f} km/h", RED)
    add_bullets(
        slide,
        [
            "Raw fields: periode_data, bulan, nama_jalan, kecepatan.",
            "The speed column uses Indonesian decimal commas and is cleaned into numeric km/h.",
            "The current dataset is useful for identifying slow corridors, but not enough for complete causal calibration.",
        ],
        0.85,
        3.15,
        11.6,
        2.4,
    )
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "Slowest Corridors", "Mean peak-hour speed by corridor, 2023")
    slide.shapes.add_picture(str(SLOWEST_FIGURE), Inches(0.75), Inches(1.35), width=Inches(7.1))
    table_df = ranking.head(6).copy()
    table_df["mean_speed_kmh"] = table_df["mean_speed_kmh"].map(lambda x: f"{x:.2f}")
    table_df = table_df[["bottleneck_rank", "corridor", "mean_speed_kmh", "observations"]]
    table_df.columns = ["Rank", "Corridor", "Mean km/h", "Obs."]
    add_table(slide, table_df, 8.15, 1.58, 4.45, 2.65)
    add_bullets(
        slide,
        [
            f"Slowest corridor: {slowest['corridor']} at {slowest['mean_speed_kmh']:.2f} km/h.",
            "These corridors are candidates for deeper bottleneck analysis.",
        ],
        8.2,
        4.55,
        4.2,
        1.3,
    )
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "Monthly Pattern", "Average speed varies, but congestion signal remains persistent")
    slide.shapes.add_picture(str(MONTHLY_FIGURE), Inches(0.75), Inches(1.35), width=Inches(7.1))
    add_metric_card(slide, 8.25, 1.55, 3.75, "lowest monthly mean", f"{slowest_month['mean_speed_kmh']:.2f} km/h", RED)
    add_metric_card(slide, 8.25, 2.85, 3.75, "month", f"{int(slowest_month['month'])}/2023", RED)
    add_metric_card(slide, 8.25, 4.15, 3.75, "highest monthly mean", f"{fastest_month['mean_speed_kmh']:.2f} km/h", GREEN)
    add_metric_card(slide, 8.25, 5.45, 3.75, "month", f"{int(fastest_month['month'])}/2023", GREEN)
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "How To Read The Congestion Mechanism", "A simple demand-capacity story")
    add_bullets(
        slide,
        [
            "Vehicle demand pushes more vehicles into limited corridor space.",
            "At peak hours, density rises and effective capacity becomes binding.",
            "Higher density reduces average speed.",
            "Local slowdowns become bottlenecks.",
            "Bottlenecks can create upstream queue-like waves, which is why traffic can remain slow even after the initial trigger changes.",
        ],
        0.9,
        1.55,
        6.0,
        4.8,
    )
    flow = [
        ("Demand", 7.25, 1.65, BLUE),
        ("Density", 9.2, 1.65, BLUE),
        ("Speed drops", 7.25, 3.25, RED),
        ("Bottleneck", 9.2, 3.25, RED),
        ("Queue waves", 8.25, 4.85, NAVY),
    ]
    for label, left, top, color in flow:
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(1.65), Inches(0.75))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        txt = shape.text_frame.paragraphs[0]
        txt.text = label
        txt.font.bold = True
        txt.font.size = Pt(13)
        txt.font.color.rgb = WHITE
        txt.alignment = PP_ALIGN.CENTER
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "Interpretation", "What the current evidence supports")
    add_bullets(
        slide,
        [
            "The corridor-speed dataset supports a bottleneck screening analysis.",
            "Slow corridors can be prioritized for follow-up using demand, geometry, public transport, and incident data.",
            "The data does not yet prove a full causal model of Jakarta congestion.",
            "A defensible next step is to connect corridor speed with vehicle demand and a simplified traffic-flow simulation.",
        ],
        0.85,
        1.55,
        11.8,
        4.4,
    )
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "Connection To Thesis Enhancement", "From thesis prototype to reproducible traffic analysis")
    add_bullets(
        slide,
        [
            "Original thesis direction: fluid-inspired traffic modelling.",
            "Repository enhancement: reproducible Python simulation and dashboard.",
            "Jakarta repository: separate empirical case study using open data.",
            "Future bridge: use Jakarta corridor speeds to calibrate scenario parameters in a simplified traffic-flow model.",
        ],
        0.85,
        1.55,
        11.7,
        4.8,
    )
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "Next Work", "Data layers needed before stronger claims")
    add_bullets(
        slide,
        [
            "Add BPS vehicle ownership and road-length context.",
            "Find corridor-level volume or count data for Jakarta.",
            "Map corridor speed to candidate bottleneck strength.",
            "Compare simulated speed-density patterns against observed corridor speeds.",
            "Document limitations explicitly before turning this into an article or paper draft.",
        ],
        0.85,
        1.55,
        11.7,
        4.8,
    )
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "Sources And Repository")
    add_bullets(
        slide,
        [
            "Satu Data Jakarta / Data.go.id: Kecepatan Rata-rata di 41 Koridor Jalan Utama pada Jam Sibuk.",
            "CSV file downloaded from the Satu Data Jakarta export endpoint.",
            "Repository: https://github.com/kandref/jakarta-traffic-congestion-analysis",
            f"Generated deck: {OUTPUT_PATH.name}",
        ],
        0.85,
        1.55,
        11.7,
        4.8,
    )
    add_footer(slide)

    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved presentation to {OUTPUT_PATH}")


if __name__ == "__main__":
    build_deck()
